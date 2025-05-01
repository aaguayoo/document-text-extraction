# Codigo para extraer texto de un archivo PPTX #
# Generado para el curso de desarrollo de Software 2025
# Created by: Alexia Medina
# Last update: 30/04/2025

import os
from pptx import Presentation

#@dataclass
class PPTXTextExtraction:
    """
    Clase para extraer texto de archivos .pptx.

    Atributos:
    - filepath: str Path del archivo PPTX del que se desea extraer el texto

    Raises:
    - ValueError: Si el archivo no es de tipo PPTX

    """

    def __init__(self):
        """
        Inicializa la clase. Solicita nombre del archivo PPTX
        Revisa que se trate de un archivo PPTX, si no, da error

        Parametros:
            -filepath: Ruta al archivo PPTX

        Raises:
            -ValueError: Si el archivo no tiene extensión .pptx

        """

        filepath_raw = input("Indique el nombre de la presentación de la que desea extraer el texto:")
        filepath = filepath_raw.strip('""')

        path_ = os.path.abspath(filepath[filepath.rfind('\\')])
        file = filepath[filepath.rfind('\\')+1:]

        if not file.lower().endswith(".pptx"):
            raise ValueError("El archivo NO es de tipo PPTX")

        self.filepath = os.path.join(path_,file)
        self.presentation = Presentation(os.path.join(path_,file))

    def extract_text(self):
        """
        Funcion que extrae el texto de cada diapositiva de una presentacion de tipo PPTX.

        Return:
            - texto_diapos: list Lista con un str que corresponde al texto de TODAS las diapositivas, separando con espacio lo de cada slide.

        """

        # Lista final que contendrá TODO el texto de TODAS las diapositivas
        texto_all_diapos = []

        # Recorro cada slide:
        for diapo in self.presentation.slides:

            texto_en_onediapo = []

            for shape in diapo.shapes:

                # Reviso si tiene atributo de texto
                if hasattr(shape, "text"):

                    texto_en_onediapo.append(shape.text.strip())

            texto_all_diapos.append(" ".join(texto_en_onediapo))

        return texto_all_diapos

    def print_text(self):
        """
        Imprime el texto de cada diapositiva, numeradas.
        """
        text_by_slide = self.extract_text()

        for i, slide_text in enumerate(text_by_slide):
            print(f"DIAPOSITIVA {i + 1}:\n{slide_text}\n")

pptx_extractor = PPTXTextExtraction()  # esto pedirá la ruta por input
pptx_extractor.print_text() # esto imprime los textos
            
